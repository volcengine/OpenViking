# Copyright (c) 2026 Beijing Volcano Engine Technology Co., Ltd.
# SPDX-License-Identifier: AGPL-3.0
"""
PowerPoint (.pptx) parser for OpenViking.

Converts PowerPoint presentations to Markdown then parses using MarkdownParser.
Inspired by microsoft/markitdown approach.
"""

import asyncio
from pathlib import Path
from typing import List, Optional, Union

from openviking.parse.base import ParseResult
from openviking.parse.parsers.base_parser import BaseParser
from openviking_cli.utils.config.parser_config import AnydocConfig, ParserConfig
from openviking_cli.utils.logger import get_logger

logger = get_logger(__name__)

_LEGACY_SUPPORTED_EXTENSIONS = frozenset({".pptx"})


class PowerPointParser(BaseParser):
    """
    PowerPoint presentation parser for OpenViking.

    Supports: .pptx, .ppt, .pptm, .pps, .ppsx, .ppsm, .pot, .odp

    Converts PowerPoint presentations to Markdown using anydoc, then delegates
    to MarkdownParser for tree structure creation. The python-pptx converter
    remains available for .pptx when anydoc is disabled or fallback is enabled.
    """

    def __init__(
        self,
        config: Optional[ParserConfig] = None,
        anydoc_config: Optional[AnydocConfig] = None,
        extract_notes: bool = False,
    ):
        """
        Initialize PowerPoint parser.

        Args:
            config: Parser configuration
            anydoc_config: Shared anydoc converter configuration
            extract_notes: Whether the legacy converter extracts speaker notes
        """
        from openviking.parse.parsers.markdown import MarkdownParser

        self._md_parser = MarkdownParser(config=config)
        self.config = config or ParserConfig()
        self.anydoc_config = anydoc_config or AnydocConfig()
        self.extract_notes = extract_notes

    @property
    def supported_extensions(self) -> List[str]:
        return [".pptx", ".ppt", ".pptm", ".pps", ".ppsx", ".ppsm", ".pot", ".odp"]

    async def parse(self, source: Union[str, Path], instruction: str = "", **kwargs) -> ParseResult:
        """Parse PowerPoint presentation from file path."""
        path = Path(source)

        if path.exists():
            from openviking_cli.utils.storage import get_storage

            storage = get_storage()
            resource_name = kwargs.get("resource_name") or kwargs.get("source_name") or path.stem
            source_format = path.suffix.lstrip(".").lower() or "pptx"

            if self.anydoc_config.enable:
                from openviking.parse.parsers.anydoc_converter import AnyDocConverter

                try:
                    conversion = await asyncio.to_thread(
                        AnyDocConverter().convert,
                        path,
                        resource_name=resource_name,
                        storage=storage,
                    )
                    markdown_content = conversion.markdown
                    source_format = conversion.source_format or source_format
                except Exception:
                    if (
                        not self.anydoc_config.fallback_to_legacy
                        or path.suffix.lower() not in _LEGACY_SUPPORTED_EXTENSIONS
                    ):
                        raise
                    logger.warning(
                        "[PowerPointParser] anydoc conversion failed for %s; "
                        "using legacy converter",
                        path.name,
                        exc_info=True,
                    )
                    markdown_content = await self._legacy_convert(path)
            else:
                if path.suffix.lower() not in _LEGACY_SUPPORTED_EXTENSIONS:
                    raise RuntimeError(
                        "anydoc conversion is disabled and no legacy converter is available "
                        f"for {path.suffix.lower() or 'this format'}"
                    )
                markdown_content = await self._legacy_convert(path)

            markdown_kwargs = dict(kwargs)
            allowed_media_dirs = list(markdown_kwargs.get("allowed_media_dirs") or [])
            if storage.media_dir not in allowed_media_dirs:
                allowed_media_dirs.append(storage.media_dir)
            markdown_kwargs.update(
                source_path=str(path),
                base_dir=path.parent,
                allowed_media_dirs=allowed_media_dirs,
            )
            result = await self._md_parser.parse_content(
                markdown_content,
                instruction=instruction,
                **markdown_kwargs,
            )
            result.source_format = source_format
        else:
            result = await self._md_parser.parse_content(
                str(source), instruction=instruction, **kwargs
            )
            result.source_format = "pptx"
        result.parser_name = "PowerPointParser"
        return result

    async def _legacy_convert(self, path: Path) -> str:
        """Run the python-pptx converter off the event loop."""
        import pptx

        return await asyncio.to_thread(self._convert_to_markdown, path, pptx)

    async def parse_content(
        self, content: str, source_path: Optional[str] = None, instruction: str = "", **kwargs
    ) -> ParseResult:
        """Parse content - delegates to MarkdownParser."""
        result = await self._md_parser.parse_content(content, source_path, **kwargs)
        result.source_format = "pptx"
        result.parser_name = "PowerPointParser"
        return result

    def _convert_to_markdown(self, path: Path, pptx) -> str:
        """Convert PowerPoint presentation to Markdown string."""
        prs = pptx.Presentation(path)
        markdown_parts = []
        slide_count = len(prs.slides)

        for idx, slide in enumerate(prs.slides, 1):
            slide_parts = []
            slide_parts.append(f"## Slide {idx}/{slide_count}")

            title = self._extract_slide_title(slide)
            if title:
                slide_parts.append(f"### {title}")

            content = self._extract_slide_content(slide)
            if content:
                slide_parts.append(content)

            if self.extract_notes and slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    slide_parts.append(f"**Notes:** {notes}")

            markdown_parts.append("\n\n".join(slide_parts))

        return "\n\n---\n\n".join(markdown_parts)

    def _extract_slide_title(self, slide) -> str:
        """Extract title from a slide."""
        from pptx.enum.shapes import PP_PLACEHOLDER

        for shape in slide.shapes:
            if shape.is_placeholder:
                ph_type = shape.placeholder_format.type
                if ph_type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
                    return shape.text.strip()
        return ""

    def _extract_slide_content(self, slide) -> str:
        """Extract content from slide shapes."""
        from pptx.enum.shapes import PP_PLACEHOLDER

        content_parts = []

        for shape in slide.shapes:
            if shape.is_placeholder:
                ph_type = shape.placeholder_format.type
                if ph_type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
                    continue

            if hasattr(shape, "text") and shape.text.strip():
                if shape.has_table:
                    content_parts.append(self._convert_table(shape.table))
                else:
                    text = shape.text.strip()
                    if text:
                        content_parts.append(text)

        return "\n\n".join(content_parts)

    def _convert_table(self, table) -> str:
        """Convert PowerPoint table to markdown format."""
        if not table.rows:
            return ""

        rows = []
        for row in table.rows:
            row_data = [cell.text.strip() for cell in row.cells]
            rows.append(row_data)

        from openviking.parse.base import format_table_to_markdown

        return format_table_to_markdown(rows, has_header=True)
