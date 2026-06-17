# Copyright (c) 2026 Beijing Volcano Engine Technology Co., Ltd.
# SPDX-License-Identifier: AGPL-3.0
"""Tests for PowerPointParser image and table extraction."""

import io
import json
import re
import struct
import zlib
from pathlib import Path

import pytest

from openviking.parse.parsers.powerpoint import PowerPointParser


def _make_png_bytes(width: int = 4, height: int = 4) -> bytes:
    """Build a tiny valid PNG without depending on Pillow."""

    def chunk(tag: bytes, data: bytes) -> bytes:
        return (
            struct.pack(">I", len(data))
            + tag
            + data
            + struct.pack(">I", zlib.crc32(tag + data) & 0xFFFFFFFF)
        )

    sig = b"\x89PNG\r\n\x1a\n"
    ihdr = chunk(b"IHDR", struct.pack(">IIBBBBB", width, height, 8, 2, 0, 0, 0))
    raw = b"".join(b"\x00" + b"\xff\x00\x00" * width for _ in range(height))
    idat = chunk(b"IDAT", zlib.compress(raw))
    iend = chunk(b"IEND", b"")
    return sig + ihdr + idat + iend


@pytest.fixture
def sample_pptx(tmp_path: Path) -> Path:
    """Build a tiny pptx with one slide containing a title, a table, and a picture."""
    pptx = pytest.importorskip("pptx")
    from pptx.util import Inches

    prs = pptx.Presentation()
    slide_layout = prs.slide_layouts[5]  # Title only
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Quarterly Results"

    # Table: 3 rows x 2 cols
    rows, cols = 3, 2
    left = top = Inches(1)
    width = Inches(4)
    height = Inches(2)
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    table.cell(0, 0).text = "Region"
    table.cell(0, 1).text = "Revenue"
    table.cell(1, 0).text = "EMEA"
    table.cell(1, 1).text = "100"
    table.cell(2, 0).text = "APAC"
    table.cell(2, 1).text = "200"

    # Picture
    image_stream = io.BytesIO(_make_png_bytes())
    slide.shapes.add_picture(image_stream, Inches(5), Inches(1), Inches(2), Inches(2))

    out = tmp_path / "sample.pptx"
    prs.save(str(out))
    return out


@pytest.fixture
def isolated_storage(tmp_path: Path, monkeypatch):
    """Pin the global StoragePath at a per-test temp dir to keep cwd clean."""
    from openviking_cli.utils import storage as storage_mod

    storage = storage_mod.StoragePath(base_path=tmp_path)
    storage.ensure_dirs()
    monkeypatch.setattr(storage_mod, "get_storage", lambda base_path=None: storage)
    return storage


@pytest.mark.asyncio
async def test_pptx_image_and_table_extraction(sample_pptx: Path, isolated_storage):
    parser = PowerPointParser()
    pptx = pytest.importorskip("pptx")

    # Pull the intermediate markdown by calling the sync helper directly.
    md = parser._convert_to_markdown(
        sample_pptx, pptx, resource_name="sample", storage=isolated_storage
    )

    # Slide heading preserved.
    assert "## Slide 1/1" in md
    assert "### Quarterly Results" in md

    # Picture saved to media dir under deterministic name.
    images_dir = isolated_storage.get_resource_media_dir("sample", "images")
    saved = list(images_dir.glob("slide1_image1.*"))
    assert saved, f"expected slide1_image1.* under {images_dir}, got {list(images_dir.iterdir())}"

    # Markdown contains a relative reference to that saved image.
    rel = saved[0].relative_to(isolated_storage.media_dir)
    assert f"]({rel})" in md

    # Fenced JSON block with the table rows (order preserved).
    json_match = re.search(r"```json\n(.*?)\n```", md, re.DOTALL)
    assert json_match, f"expected fenced json block in markdown:\n{md}"
    payload = json.loads(json_match.group(1))
    assert payload == {
        "rows": [
            ["Region", "Revenue"],
            ["EMEA", "100"],
            ["APAC", "200"],
        ]
    }

    # Plain markdown table emitted after the JSON block (additive).
    after_json = md[json_match.end() :]
    assert re.search(r"\|\s*Region\s*\|\s*Revenue\s*\|", after_json)
    assert re.search(r"\|\s*-+\s*\|\s*-+\s*\|", after_json)
    assert re.search(r"\|\s*EMEA\s*\|\s*100\s*\|", after_json)


@pytest.mark.asyncio
async def test_pptx_image_filenames_are_deterministic(sample_pptx: Path, isolated_storage):
    """Re-parsing the same pptx must overwrite, not pile up duplicates."""
    parser = PowerPointParser()
    pptx = pytest.importorskip("pptx")

    parser._convert_to_markdown(sample_pptx, pptx, resource_name="sample", storage=isolated_storage)
    parser._convert_to_markdown(sample_pptx, pptx, resource_name="sample", storage=isolated_storage)

    images_dir = isolated_storage.get_resource_media_dir("sample", "images")
    saved = sorted(images_dir.glob("slide1_image*.*"))
    # Same deterministic name reused; second pass must not produce slide1_image2.
    assert len(saved) == 1, f"expected one image, got {[p.name for p in saved]}"
