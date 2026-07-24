# Copyright (c) 2026 Beijing Volcano Engine Technology Co., Ltd.
# SPDX-License-Identifier: AGPL-3.0
"""PowerPoint embedded-image extraction regressions."""

import hashlib
from pathlib import Path
from types import SimpleNamespace
from unittest.mock import AsyncMock, MagicMock

import pptx
import pytest
from PIL import Image
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER

from openviking.parse.parsers.powerpoint import PowerPointParser
from openviking_cli.utils.storage import StoragePath


def _write_png(path: Path, color: tuple[int, int, int]) -> bytes:
    Image.new("RGB", (32, 24), color=color).save(path)
    return path.read_bytes()


def _picture_placeholder_layout(prs):
    for layout in prs.slide_layouts:
        if any(
            placeholder.placeholder_format.type == PP_PLACEHOLDER.PICTURE
            for placeholder in layout.placeholders
        ):
            return layout
    raise AssertionError("default PowerPoint template has no picture placeholder")


def test_convert_to_markdown_persists_picture_shapes_in_slide_order(tmp_path: Path):
    parser = PowerPointParser()
    media_dir = tmp_path / "media"
    digest = hashlib.sha256(b"embedded-png").hexdigest()[:12]
    filename = f"image1_{digest}"
    resource_images = media_dir / "deck" / "images"
    saved_path = resource_images / f"{filename}.png"
    storage = SimpleNamespace(
        media_dir=media_dir,
        get_resource_media_dir=MagicMock(return_value=resource_images),
        save_image=MagicMock(return_value=saved_path),
    )

    picture = SimpleNamespace(
        is_placeholder=False,
        shape_type=MSO_SHAPE_TYPE.PICTURE,
        image=SimpleNamespace(blob=b"embedded-png", ext="png"),
    )
    text = SimpleNamespace(
        is_placeholder=False,
        shape_type=MSO_SHAPE_TYPE.TEXT_BOX,
        text="Caption after image",
        has_table=False,
    )
    presentation = SimpleNamespace(
        slides=[SimpleNamespace(shapes=[picture, text], has_notes_slide=False)]
    )
    pptx_module = SimpleNamespace(Presentation=lambda _path: presentation)

    markdown = parser._convert_to_markdown(
        tmp_path / "deck.pptx", pptx_module, resource_name="deck", storage=storage
    )

    storage.save_image.assert_called_once_with(
        "deck", b"embedded-png", filename=filename, extension=".png"
    )
    assert f"![{filename}](images/{filename}.png)" in markdown
    assert markdown.index(f"![{filename}]") < markdown.index("Caption after image")


def test_convert_to_markdown_recurses_into_grouped_pictures(tmp_path: Path):
    parser = PowerPointParser()
    media_dir = tmp_path / "media"
    digest = hashlib.sha256(b"grouped-png").hexdigest()[:12]
    filename = f"image1_{digest}"
    resource_images = media_dir / "deck" / "images"
    saved_path = resource_images / f"{filename}.png"
    storage = SimpleNamespace(
        media_dir=media_dir,
        get_resource_media_dir=MagicMock(return_value=resource_images),
        save_image=MagicMock(return_value=saved_path),
    )
    picture = SimpleNamespace(
        is_placeholder=False,
        shape_type=MSO_SHAPE_TYPE.PICTURE,
        image=SimpleNamespace(blob=b"grouped-png", ext="png"),
    )
    group = SimpleNamespace(
        is_placeholder=False,
        shape_type=MSO_SHAPE_TYPE.GROUP,
        shapes=[picture],
    )
    slide = SimpleNamespace(shapes=[group], has_notes_slide=False)

    markdown = parser._extract_slide_content(
        slide, resource_name="deck", storage=storage, image_counter=[0]
    )

    assert markdown == f"![{filename}](images/{filename}.png)"


def test_real_picture_placeholder_is_extracted(tmp_path: Path):
    parser = PowerPointParser()
    storage = StoragePath(base_path=tmp_path)
    image_path = tmp_path / "placeholder.png"
    image_bytes = _write_png(image_path, (220, 40, 40))
    presentation_path = tmp_path / "placeholder.pptx"

    prs = pptx.Presentation()
    slide = prs.slides.add_slide(_picture_placeholder_layout(prs))
    picture_placeholder = next(
        placeholder
        for placeholder in slide.placeholders
        if placeholder.placeholder_format.type == PP_PLACEHOLDER.PICTURE
    )
    picture_placeholder.insert_picture(str(image_path))
    prs.save(presentation_path)

    markdown = parser._convert_to_markdown(
        presentation_path, pptx, resource_name="placeholder", storage=storage
    )

    digest = hashlib.sha256(image_bytes).hexdigest()[:12]
    assert f"![image1_{digest}]" in markdown
    assert (storage.media_dir / "placeholder" / "images" / f"image1_{digest}.png").exists()


def test_same_stem_presentations_do_not_overwrite_extracted_images(tmp_path: Path):
    parser = PowerPointParser()
    storage = StoragePath(base_path=tmp_path)
    markdown_outputs = []
    expected_blobs = []

    for directory, color in (("first", (20, 40, 220)), ("second", (40, 180, 60))):
        source_dir = tmp_path / directory
        source_dir.mkdir()
        image_path = source_dir / "picture.png"
        expected_blobs.append(_write_png(image_path, color))
        presentation_path = source_dir / "deck.pptx"
        prs = pptx.Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(str(image_path), 0, 0)
        prs.save(presentation_path)
        markdown_outputs.append(
            parser._convert_to_markdown(
                presentation_path, pptx, resource_name="deck", storage=storage
            )
        )

    stored_images = sorted((storage.media_dir / "deck" / "images").glob("image1_*.png"))
    assert len(stored_images) == 2
    assert {path.read_bytes() for path in stored_images} == set(expected_blobs)
    assert markdown_outputs[0] != markdown_outputs[1]


def test_convert_picture_failure_is_non_fatal(tmp_path: Path):
    parser = PowerPointParser()
    storage = SimpleNamespace(
        media_dir=tmp_path,
        save_image=MagicMock(side_effect=OSError("media unavailable")),
    )
    picture = SimpleNamespace(image=SimpleNamespace(blob=b"data", ext="jpeg"))

    assert parser._convert_picture(picture, "deck", storage, [0]) == ""


@pytest.mark.asyncio
async def test_parse_does_not_authorize_sibling_resource_media(monkeypatch, tmp_path: Path):
    from openviking_cli.utils import storage as storage_module

    parser = PowerPointParser()
    storage = StoragePath(base_path=tmp_path / "storage")
    sibling_image = storage.get_resource_media_dir("sibling", "images") / "private.png"
    sibling_image.write_bytes(b"private sibling bytes")
    source = tmp_path / "incoming" / "deck.pptx"
    source.parent.mkdir()
    source.write_bytes(b"placeholder")

    monkeypatch.setattr(storage_module, "get_storage", lambda: storage)
    monkeypatch.setattr(
        parser,
        "_convert_to_markdown",
        lambda *args, **kwargs: "![x](sibling/images/private.png)",
    )
    parser._md_parser.parse_content = AsyncMock(
        return_value=SimpleNamespace(source_format="markdown", parser_name="MarkdownParser")
    )

    await parser.parse(source, resource_name="deck")

    parse_kwargs = parser._md_parser.parse_content.await_args.kwargs
    allowed_media_dirs = parse_kwargs["allowed_media_dirs"]
    assert storage.media_dir not in allowed_media_dirs
    assert allowed_media_dirs == []
    assert parse_kwargs["preferred_media_paths"] == {}
    assert (
        parser._md_parser._resolve_image_path(
            "sibling/images/private.png",
            parse_kwargs["base_dir"],
            allowed_media_dirs,
            parse_kwargs["preferred_media_paths"],
        )
        is None
    )


@pytest.mark.asyncio
async def test_generated_media_mapping_wins_source_and_caller_shadows_and_is_cleaned(
    monkeypatch, tmp_path: Path
):
    from openviking_cli.utils import storage as storage_module

    parser = PowerPointParser()
    storage = StoragePath(base_path=tmp_path / "storage")
    source_dir = tmp_path / "incoming"
    source_dir.mkdir()
    image_path = source_dir / "picture.png"
    image_bytes = _write_png(image_path, (25, 100, 220))
    source = source_dir / "Deck (Final).pptx"
    prs = pptx.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(str(image_path), 0, 0)
    prs.save(source)

    digest = hashlib.sha256(image_bytes).hexdigest()[:12]
    relative_image = Path("images") / f"image1_{digest}.png"
    caller_root = tmp_path / "caller-media"
    caller_shadow = caller_root / relative_image
    caller_shadow.parent.mkdir(parents=True)
    caller_shadow.write_bytes(b"caller shadow")
    source_shadow = source_dir / relative_image
    source_shadow.parent.mkdir(parents=True)
    source_shadow.write_bytes(b"source shadow")

    monkeypatch.setattr(storage_module, "get_storage", lambda: storage)
    observed = {}

    async def capture_parse(content, *args, **kwargs):
        preferred_media_paths = kwargs["preferred_media_paths"]
        generated_image = preferred_media_paths[relative_image.as_posix()]
        resource_root = generated_image.parent.parent
        table_image = resource_root / "tables" / "private.png"
        table_image.parent.mkdir()
        table_image.write_bytes(b"private table bytes")
        observed.update(
            content=content,
            allowed_media_dirs=kwargs["allowed_media_dirs"],
            preferred_media_paths=dict(preferred_media_paths),
            generated_image=generated_image,
            resource_root=resource_root,
            generated_bytes=generated_image.read_bytes(),
            resolved=parser._md_parser._resolve_image_path(
                relative_image.as_posix(),
                kwargs["base_dir"],
                kwargs["allowed_media_dirs"],
                preferred_media_paths,
            ),
            cross_media=parser._md_parser._resolve_image_path(
                "tables/private.png",
                kwargs["base_dir"],
                kwargs["allowed_media_dirs"],
                preferred_media_paths,
            ),
        )
        return SimpleNamespace(source_format="markdown", parser_name="MarkdownParser")

    parser._md_parser.parse_content = AsyncMock(side_effect=capture_parse)

    await parser.parse(
        source,
        resource_name="Deck (Final)",
        allowed_media_dirs=[caller_root],
    )

    assert f"![image1_{digest}]({relative_image.as_posix()})" in observed["content"]
    assert "Deck (Final)" not in observed["content"]
    assert observed["allowed_media_dirs"] == [caller_root]
    assert list(observed["preferred_media_paths"]) == [relative_image.as_posix()]
    assert observed["generated_bytes"] == image_bytes
    assert observed["resolved"] == observed["generated_image"].resolve()
    assert observed["resolved"] != source_shadow.resolve()
    assert observed["resolved"] != caller_shadow.resolve()
    assert observed["cross_media"] is None
    assert not observed["resource_root"].exists()


async def _capture_generated_path(
    parser: PowerPointParser,
    source: Path,
    *,
    resource_name: str,
) -> Path:
    captured = {}

    async def capture_parse(_content, *args, **kwargs):
        generated_paths = kwargs["preferred_media_paths"]
        assert len(generated_paths) == 1
        captured["path"] = next(iter(generated_paths.values()))
        assert captured["path"].exists()
        return SimpleNamespace(source_format="markdown", parser_name="MarkdownParser")

    parser._md_parser.parse_content = AsyncMock(side_effect=capture_parse)
    await parser.parse(source, resource_name=resource_name)
    assert not captured["path"].parent.parent.exists()
    return captured["path"]


@pytest.mark.asyncio
async def test_same_stem_presentations_use_distinct_parse_namespaces(monkeypatch, tmp_path: Path):
    from openviking_cli.utils import storage as storage_module

    parser = PowerPointParser()
    storage = StoragePath(base_path=tmp_path / "storage")
    monkeypatch.setattr(storage_module, "get_storage", lambda: storage)
    sources = []
    for directory in ("first", "second"):
        source_dir = tmp_path / directory
        source_dir.mkdir()
        image = source_dir / "picture.png"
        _write_png(image, (20, 40, 220))
        source = source_dir / "deck.pptx"
        prs = pptx.Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(str(image), 0, 0)
        prs.save(source)
        sources.append(source)

    paths = [
        await _capture_generated_path(parser, source, resource_name="deck") for source in sources
    ]

    assert paths[0].parent.parent != paths[1].parent.parent


@pytest.mark.asyncio
async def test_sanitized_resource_name_collisions_use_distinct_parse_namespaces(
    monkeypatch, tmp_path: Path
):
    from openviking_cli.utils import storage as storage_module

    parser = PowerPointParser()
    storage = StoragePath(base_path=tmp_path / "storage")
    monkeypatch.setattr(storage_module, "get_storage", lambda: storage)
    image = tmp_path / "picture.png"
    _write_png(image, (40, 180, 60))
    source = tmp_path / "deck.pptx"
    prs = pptx.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(str(image), 0, 0)
    prs.save(source)

    paths = [
        await _capture_generated_path(parser, source, resource_name=name)
        for name in ("deck/a", "deck:a")
    ]

    assert StoragePath._sanitize_name("deck/a") == StoragePath._sanitize_name("deck:a")
    assert paths[0].parent.parent != paths[1].parent.parent
